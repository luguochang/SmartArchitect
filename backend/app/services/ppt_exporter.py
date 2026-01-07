"""
PPT Exporter Service
Creates PowerPoint presentations from architecture diagrams
"""

import logging
from io import BytesIO
from typing import List
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    Presentation = None

from app.models.schemas import Node, Edge

logger = logging.getLogger(__name__)


class PPTExporter:
    """Service for exporting architecture diagrams to PowerPoint"""

    def __init__(self):
        if Presentation is None:
            raise ImportError(
                "python-pptx is not installed. Install with: pip install python-pptx"
            )

    def create_presentation(
        self, nodes: List[Node], edges: List[Edge], title: str = "Architecture Diagram"
    ) -> bytes:
        """Create a PowerPoint presentation from nodes and edges"""

        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Slide 1: Title slide
            self._add_title_slide(prs, title)

            # Slide 2: Architecture diagram
            self._add_diagram_slide(prs, nodes, edges, title)

            # Slide 3: Component breakdown
            self._add_components_slide(prs, nodes)

            # Slide 4: Connections breakdown
            self._add_connections_slide(prs, nodes, edges)

            # Save to bytes
            ppt_bytes = BytesIO()
            prs.save(ppt_bytes)
            ppt_bytes.seek(0)

            logger.info(f"Created PowerPoint presentation with {len(prs.slides)} slides")
            return ppt_bytes.getvalue()

        except Exception as e:
            logger.error(f"Failed to create PowerPoint: {e}", exc_info=True)
            raise

    def _add_title_slide(self, prs: Presentation, title: str):
        """Add title slide"""

        slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}\nby SmartArchitect AI"

    def _add_diagram_slide(self, prs: Presentation, nodes: List[Node], edges: List[Edge], title: str):
        """Add architecture diagram slide"""

        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_shape.text_frame
        title_frame.text = f"{title} - Architecture Diagram"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        # Normalize node positions to fit slide
        if len(nodes) > 0:
            positions = [(node.position.x, node.position.y) for node in nodes]
            min_x = min(p[0] for p in positions)
            max_x = max(p[0] for p in positions)
            min_y = min(p[1] for p in positions)
            max_y = max(p[1] for p in positions)

            # Available area: 9 inches wide, 6 inches tall (leave margins)
            width_range = max_x - min_x if max_x > min_x else 1
            height_range = max_y - min_y if max_y > min_y else 1

            scale_x = 8 / width_range
            scale_y = 5 / height_range
            scale = min(scale_x, scale_y) * 0.8  # Add some padding

            # Center the diagram
            offset_x = Inches(1)
            offset_y = Inches(2)

            # Draw edges first (so they appear behind nodes)
            for edge in edges:
                source_node = next((n for n in nodes if n.id == edge.source), None)
                target_node = next((n for n in nodes if n.id == edge.target), None)

                if source_node and target_node:
                    x1 = offset_x + Inches((source_node.position.x - min_x) * scale / 100)
                    y1 = offset_y + Inches((source_node.position.y - min_y) * scale / 100)
                    x2 = offset_x + Inches((target_node.position.x - min_x) * scale / 100)
                    y2 = offset_y + Inches((target_node.position.y - min_y) * scale / 100)

                    # Draw connector line
                    connector = slide.shapes.add_connector(
                        1,  # msoConnectorStraight
                        x1 + Inches(0.75), y1 + Inches(0.25),  # Start point (center of node)
                        x2 + Inches(0.75), y2 + Inches(0.25)   # End point (center of node)
                    )
                    connector.line.color.rgb = RGBColor(128, 128, 128)
                    connector.line.width = Pt(2)

            # Draw nodes
            for node in nodes:
                x = offset_x + Inches((node.position.x - min_x) * scale / 100)
                y = offset_y + Inches((node.position.y - min_y) * scale / 100)

                # Determine node color based on type
                color = self._get_node_color(node.type or "default")

                # Add node shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(1.5), Inches(0.5)
                )

                # Style the shape
                shape.fill.solid()
                shape.fill.fore_color.rgb = color["fill"]
                shape.line.color.rgb = color["border"]
                shape.line.width = Pt(2)

                # Add text
                text_frame = shape.text_frame
                text_frame.text = node.data.label
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 0, 0)

    def _add_components_slide(self, prs: Presentation, nodes: List[Node]):
        """Add component breakdown slide"""

        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "System Components"

        # Add content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = f"Total Components: {len(nodes)}"

        # Group by type
        types = {}
        for node in nodes:
            node_type = node.type or "default"
            if node_type not in types:
                types[node_type] = []
            types[node_type].append(node.data.label)

        for node_type, labels in types.items():
            p = tf.add_paragraph()
            p.text = f"\n{node_type.upper()}: {', '.join(labels)}"
            p.level = 0
            p.font.size = Pt(14)

    def _add_connections_slide(self, prs: Presentation, nodes: List[Node], edges: List[Edge]):
        """Add connections breakdown slide"""

        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Component Connections"

        # Add content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = f"Total Connections: {len(edges)}"

        for edge in edges:
            source_node = next((n for n in nodes if n.id == edge.source), None)
            target_node = next((n for n in nodes if n.id == edge.target), None)

            if source_node and target_node:
                p = tf.add_paragraph()
                label_text = f" ({edge.label})" if edge.label else ""
                p.text = f"{source_node.data.label} → {target_node.data.label}{label_text}"
                p.level = 0
                p.font.size = Pt(14)

    def _get_node_color(self, node_type: str) -> dict:
        """Get color scheme for node type"""

        colors = {
            "api": {
                "fill": RGBColor(219, 234, 254),  # Light blue
                "border": RGBColor(59, 130, 246),  # Blue
            },
            "service": {
                "fill": RGBColor(250, 245, 255),  # Light purple
                "border": RGBColor(168, 85, 247),  # Purple
            },
            "database": {
                "fill": RGBColor(236, 253, 245),  # Light green
                "border": RGBColor(16, 185, 129),  # Green
            },
            "gateway": {
                "fill": RGBColor(255, 251, 235),  # Light amber
                "border": RGBColor(245, 158, 11),  # Amber
            },
            "default": {
                "fill": RGBColor(249, 250, 251),  # Light gray
                "border": RGBColor(107, 114, 128),  # Gray
            },
        }

        return colors.get(node_type, colors["default"])


# Helper function to create exporter instance
def create_ppt_exporter() -> PPTExporter:
    """Create a PPTExporter instance"""
    return PPTExporter()
